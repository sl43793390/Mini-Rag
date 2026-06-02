import os
import logging
from typing import List

from langchain_core.documents import Document

from config import SUPPORTED_EXTENSIONS

logger = logging.getLogger(__name__)


class DocumentLoader:
    SUPPORTED_EXTENSIONS = SUPPORTED_EXTENSIONS

    @staticmethod
    def load_file(file_path: str) -> List[Document]:
        ext = os.path.splitext(file_path)[1].lower()
        file_type = DocumentLoader.SUPPORTED_EXTENSIONS.get(ext)

        if not file_type:
            logger.warning(f"Unsupported file type: {ext} for {file_path}")
            return []

        loaders = {
            "pdf": DocumentLoader._load_pdf,
            "docx": DocumentLoader._load_docx,
            "pptx": DocumentLoader._load_pptx,
            "xlsx": DocumentLoader._load_xlsx,
            "csv": DocumentLoader._load_csv,
            "text": DocumentLoader._load_text,
            "markdown": DocumentLoader._load_markdown,
        }

        loader_fn = loaders.get(file_type)
        if not loader_fn:
            logger.warning(f"No loader for file type: {file_type}")
            return []

        try:
            docs = loader_fn(file_path)
            for doc in docs:
                doc.metadata["source"] = file_path
                doc.metadata["file_type"] = file_type
                doc.metadata["file_name"] = os.path.basename(file_path)
            return docs
        except Exception as e:
            logger.error(f"Error loading {file_path}: {e}")
            return []

    @staticmethod
    def load_directory(dir_path: str) -> List[Document]:
        all_docs = []
        for root, _, files in os.walk(dir_path):
            for file_name in files:
                file_path = os.path.join(root, file_name)
                ext = os.path.splitext(file_name)[1].lower()
                if ext in DocumentLoader.SUPPORTED_EXTENSIONS:
                    docs = DocumentLoader.load_file(file_path)
                    all_docs.extend(docs)
                    logger.info(f"Loaded {len(docs)} documents from {file_name}")
        return all_docs

    @staticmethod
    def _load_pdf(file_path: str) -> List[Document]:
        from langchain_community.document_loaders import PyPDFLoader
        loader = PyPDFLoader(file_path)
        return loader.load()

    @staticmethod
    def _load_docx(file_path: str) -> List[Document]:
        from langchain_community.document_loaders import Docx2txtLoader
        loader = Docx2txtLoader(file_path)
        return loader.load()

    @staticmethod
    def _load_pptx(file_path: str) -> List[Document]:
        from pptx import Presentation
        docs = []
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                docs.append(Document(
                    page_content="\n".join(texts),
                    metadata={"page": i + 1}
                ))
        return docs

    @staticmethod
    def _load_xlsx(file_path: str) -> List[Document]:
        from langchain_community.document_loaders import UnstructuredExcelLoader
        loader = UnstructuredExcelLoader(file_path)
        return loader.load()

    @staticmethod
    def _load_csv(file_path: str) -> List[Document]:
        from langchain_community.document_loaders import CSVLoader
        loader = CSVLoader(file_path)
        return loader.load()

    @staticmethod
    def _load_text(file_path: str) -> List[Document]:
        from langchain_community.document_loaders import TextLoader
        loader = TextLoader(file_path, encoding="utf-8")
        return loader.load()

    @staticmethod
    def _load_markdown(file_path: str) -> List[Document]:
        from langchain_community.document_loaders import UnstructuredMarkdownLoader
        loader = UnstructuredMarkdownLoader(file_path)
        return loader.load()
